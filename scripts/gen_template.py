# -*- coding: utf-8 -*-
"""
生成《样本课程输入清单》可编辑文档：
  - docs/Sample_Input_Checklist.docx  (Word，空白模板可填写 + 情境领导示范)
  - docs/Sample_Input_Checklist.pptx  (PPT，逐条清单 + 示范)
单源，避免漂移。
"""
import os
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt as PPt, Inches as PIn
from pptx.dml.color import RGBColor as PRGB

OUT = os.path.join(os.path.dirname(__file__), "..", "docs")
os.makedirs(OUT, exist_ok=True)

# ---------- 单源内容 ----------
TITLE = "样本课程输入清单（课程 → AVG 转换）"
SUBTITLE = "请附上：授课 PPT（.pptx）+ 教材（PDF / Word）。填好后交给设计 Agent，先做一份 Learning Gameplay Map 供你确认。"

# 空白模板 10 项: (编号, 项目, 优先级, 填写说明, 示范填写)
ITEMS = [
    ("1", "课程基本信息", "🔴 必填",
     "课程名称、总时长、模块/章节划分、目标学员层级。并注明：你打算先切哪个模块做切片（建议挑 1 个最“有戏”的模块）。"),
    ("2", "学习目标与能力模型", "🔴 必填",
     "学完这门课，学员应该“能做 / 能判断什么”？尽量用可观测的行为动词（能识别 / 能制定 / 能调解 / 能诊断…）。"),
    ("3", "学员对象与真实场景背景", "🔴 必填",
     "谁在上这门课（岗位 / 层级 / 痛点）？他们日常最典型、最痛的真实工作情境是什么？决定角色人设与剧情是否接地气。"),
    ("4", "现有考核与课堂活动", "🔴 必填",
     "目前怎么考？有没有小组讨论 / 角色扮演 / 案例分析？这些是直接可转成 Decision（抉择）/ Reflection（复盘）节点的金矿。"),
    ("5", "标志性案例 / 真实两难", "🟠 建议",
     "课程里最出彩的那个真实故事、最纠结的抉择。这是 AVG 最好的“剧本种子”（叙事脊梁）。"),
    ("6", "“好课”参照（引导话术 / 顿悟点）", "🟠 建议",
     "你心目中“好课”长什么样：有没有引导话术、DM 脚本、或“学员顿悟的那一刻”？直接影响 Reflection 设计。"),
    ("7", "品牌 / 视觉 / 脱敏要求", "🟡 包装",
     "能否用真实公司名？要不要统一色调 / 字体？脱敏与合规要求（不外泄真实姓名等）。"),
    ("8", "期望产出形态", "🟡 包装",
     "纯剧本杀式 / 轻量分支 / 两者混合？决定离线包的组件构成与 DM 工作量。"),
    ("9", "试点成功标准", "🟡 验收",
     "这份样本做到什么程度算“成了”（能试玩？能打印？能真用于一期培训？）。"),
    ("10", "文件清单（附 PPT + 教材）", "🔴 必填",
     "列出你附上的文件：授课 PPT 路径、教材（PDF/Word）路径，并简述每份文件讲什么。"),
    ("11", "分层与前置判定（Tier × Track）", "🟠 建议",
     "这门课打算做成哪一层？——L1 轻量切片 3h / L2 标准单元 7h / L3 深度系列 14h（AVG 探索固定占 50%）。"
     "哪一类？——C1 通用改编（已有课程转译）/ C2 定制业务赋能（从业务目标萃取构建）。"
     "给出非 AVG 50% 内 导入/复盘/应用/IDP 的大致分配。Intake Agent 会复核此判定。"),
]

EXAMPLE = [
    ("1 · 课程基本信息",
     "《情境领导®（Situational Leadership®）》企业内训；1 天（6.5h）；模块：①领导风格认知 ②成熟度诊断 ③风格匹配 ④授权与跟进；"
     "目标学员：新任 / 中层经理。先切模块③“风格匹配”做切片（最有戏、最易成 Decision 节点）。"),
    ("2 · 学习目标与能力模型",
     "学完学员能——(a) 判定下属在“任务 / 关系”维度的成熟度（R1–R4）；(b) 在 S1–S4 四种领导风格中匹配正确风格；"
     "(c) 在下属成熟度变化时动态调整。行为动词：能诊断、能匹配、能调整。"),
    ("3 · 学员对象与真实场景背景",
     "新任经理（带 3–8 人团队，0–2 年管理经验）。真实情境：下属突然业绩下滑 / 资深下属不服管 / 新人频繁求助却不敢放手。"),
    ("4 · 现有考核与课堂活动",
     "现用“风格-成熟度匹配测验”+ 角色扮演（给一个下属画像，让学员选风格并说明理由）；小组讨论“最难的带人场景”。"
     "→ 角色扮演可直接转 Decision 节点，小组讨论转 Reflection。"),
    ("5 · 标志性案例 / 真实两难",
     "“老周的两难”——技术骨干老周被提拔成主管，仍习惯自己上手，团队怨声载道；上级要求他“放权”但项目紧急。"
     "这是真实两难（S1 指令式 vs S3 参与式拉扯），绝佳剧本种子。"),
    ("6 · “好课”参照",
     "顿悟时刻是“原来不是下属不行，是我用错了风格”。引导话术示例：“如果你是这个主管，此刻你最怕什么？”"),
    ("7 · 品牌 / 视觉 / 脱敏",
     "可用客户内部案例（脱敏），主色沿用客户品牌蓝；不外泄真实姓名；物料印刷需 CMYK 合规。"),
    ("8 · 期望产出形态",
     "剧本杀式（主管扮演 + 下属 NPC + 抉择）+ 轻量分支复盘。"),
    ("9 · 试点成功标准",
     "能打印成离线包、由客户内部讲师带 1 期 20 人工作坊、学员产出“我的风格调整行动计划”。"),
    ("10 · 文件清单",
     "授课 PPT：SL_PPT.pptx（含四模块、匹配测验）；教材：SL_Manual.pdf（情境领导原著节选 + 案例）。"),
    ("11 · 分层与前置判定",
     "Tier = L2 标准单元（7h：AVG 3.5h + 非 AVG 3.5h）；Track = C1 通用改编（情境领导已有完整课件）。"
     "非 AVG 50% 分配：导入 0.9h / 复盘 1.05h / 应用 0.9h / IDP 0.7h。"),
]

PIPELINE = (
    "【第 0 步 课程准入与分层判定 Intake & Tiering】→ 课程 → 能力(Capability) → 行为(Behavior) → "
    "挑战(Challenge) → 探索(Exploration) → 抉择(Decision) → 后果(Consequence) → 成长(Growth) → "
    "结局(Ending) → 复盘(Reflection) → 【应用迁移桥 Application · 强制】"
)

# ---------- CJK 字体 helper ----------
def set_cjk(run, name="Microsoft YaHei"):
    run.font.name = name
    r = run._element
    rPr = r.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = r.makeelement(qn('w:rFonts'), {})
        rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), name)

# =================== DOCX ===================
def build_docx(path):
    doc = Document()
    style = doc.styles['Normal']
    style.font.size = Pt(10.5)
    set_cjk(style.font.element.getparent().getchildren()[0] if False else style.font._element.get_or_add_rPr().getparent().getchildren()[0]) if False else None
    # title
    t = doc.add_heading(TITLE, level=0)
    for r in t.runs:
        set_cjk(r)
    p = doc.add_paragraph(SUBTITLE)
    for r in p.runs:
        set_cjk(r)
        r.font.size = Pt(10)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Part 1: 空白模板
    h1 = doc.add_heading("第一部分 · 空白模板（请逐项填写）", level=1)
    for r in h1.runs: set_cjk(r)
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    hdr = tbl.rows[0].cells
    for i, htext in enumerate(["编号","项目","优先级","填写说明（请按此提供）","你的填写"]):
        hdr[i].text = ""
        para = hdr[i].paragraphs[0]
        run = para.add_run(htext)
        run.bold = True
        set_cjk(run)
        run.font.size = Pt(9.5)
    for (num, name, pri, desc) in ITEMS:
        row = tbl.add_row().cells
        vals = [num, name, pri, desc, ""]
        for i, v in enumerate(vals):
            row[i].text = ""
            para = row[i].paragraphs[0]
            run = para.add_run(v)
            set_cjk(run)
            run.font.size = Pt(9.5)
            if i == 1:
                run.bold = True

    # Part 2: 示范
    h2 = doc.add_heading("第二部分 · 情境领导课程 · 示范填写", level=1)
    for r in h2.runs: set_cjk(r)
    for (head, body) in EXAMPLE:
        ph = doc.add_heading(head, level=2)
        for r in ph.runs: set_cjk(r)
        pb = doc.add_paragraph(body)
        for r in pb.runs: set_cjk(r)

    # Pipeline 映射
    h3 = doc.add_heading("附：课程 → Pipeline 映射示意（设计 Agent 用）", level=2)
    for r in h3.runs: set_cjk(r)
    pp = doc.add_paragraph(PIPELINE)
    for r in pp.runs: set_cjk(r)
    pp2 = doc.add_paragraph(
        "说明：清单第 2–5 项（学习目标 / 学员场景 / 现有活动 / 标志性案例）是映射的输入；"
        "第 4 项的角色扮演→Decision，小组讨论→Reflection；第 5 项的真实两难→核心抉择种子。"
    )
    for r in pp2.runs: set_cjk(r); r.font.size = Pt(9.5); r.font.color.rgb = RGBColor(0x55,0x55,0x55)
    doc.save(path)

# =================== PPTX ===================
def build_pptx(path):
    prs = Presentation()
    prs.slide_width = PIn(13.333)
    prs.slide_height = PIn(7.5)
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(PIn(0.8), PIn(2.4), PIn(11.7), PIn(1.4))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = TITLE; r.font.size = PPt(32); r.font.bold = True; r.font.color.rgb = PRGB(0x0E,0x6B,0x50)
    sb = s.shapes.add_textbox(PIn(0.8), PIn(3.9), PIn(11.7), PIn(2.0))
    stf = sb.text_frame; stf.word_wrap = True
    sr = stf.paragraphs[0].add_run(); sr.text = SUBTITLE; sr.font.size = PPt(14); sr.font.color.rgb = PRGB(0x55,0x55,0x55)

    # blank template slides
    for (num, name, pri, desc) in ITEMS:
        s = prs.slides.add_slide(blank)
        # header bar
        hb = s.shapes.add_textbox(PIn(0.6), PIn(0.5), PIn(12.1), PIn(1.0))
        htf = hb.text_frame; htf.word_wrap = True
        hr = htf.paragraphs[0].add_run()
        hr.text = f"{num}. {name}   〔{pri}〕"
        hr.font.size = PPt(24); hr.font.bold = True; hr.font.color.rgb = PRGB(0x0E,0x6B,0x50)
        # desc
        db = s.shapes.add_textbox(PIn(0.6), PIn(1.7), PIn(12.1), PIn(2.0))
        dtf = db.text_frame; dtf.word_wrap = True
        dr = dtf.paragraphs[0].add_run(); dr.text = "填写说明：\n" + desc; dr.font.size = PPt(15)
        # fill box
        fb = s.shapes.add_textbox(PIn(0.6), PIn(3.9), PIn(12.1), PIn(2.8))
        ftf = fb.text_frame; ftf.word_wrap = True
        fr = ftf.paragraphs[0].add_run(); fr.text = "你的填写：\n_________________________________________________\n\n\n\n"; fr.font.size = PPt(15)

    # example slides
    es = prs.slides.add_slide(blank)
    eb = es.shapes.add_textbox(PIn(0.6), PIn(0.5), PIn(12.1), PIn(0.8))
    ebr = eb.text_frame.paragraphs[0].add_run(); ebr.text = "情境领导课程 · 示范填写"; ebr.font.size = PPt(24); ebr.font.bold = True; ebr.font.color.rgb = PRGB(0x0E,0x6B,0x50)
    for (head, body) in EXAMPLE:
        s = prs.slides.add_slide(blank)
        hb = s.shapes.add_textbox(PIn(0.6), PIn(0.5), PIn(12.1), PIn(0.9))
        htf = hb.text_frame; htf.word_wrap = True
        hr = htf.paragraphs[0].add_run(); hr.text = head; hr.font.size = PPt(20); hr.font.bold = True; hr.font.color.rgb = PRGB(0x0E,0x6B,0x50)
        bb = s.shapes.add_textbox(PIn(0.6), PIn(1.6), PIn(12.1), PIn(5.0))
        btf = bb.text_frame; btf.word_wrap = True
        br = btf.paragraphs[0].add_run(); br.text = body; br.font.size = PPt(15)

    # pipeline slide
    ps = prs.slides.add_slide(blank)
    pb = ps.shapes.add_textbox(PIn(0.6), PIn(0.5), PIn(12.1), PIn(0.8))
    pbr = pb.text_frame.paragraphs[0].add_run(); pbr.text = "课程 → Pipeline 映射示意"; pbr.font.size = PPt(24); pbr.font.bold = True; pbr.font.color.rgb = PRGB(0x0E,0x6B,0x50)
    pbb = ps.shapes.add_textbox(PIn(0.6), PIn(1.5), PIn(12.1), PIn(5.5))
    pbtf = pbb.text_frame; pbtf.word_wrap = True
    pbr2 = pbtf.paragraphs[0].add_run(); pbr2.text = PIPELINE; pbr2.font.size = PPt(16)

    prs.save(path)

if __name__ == "__main__":
    d = os.path.abspath(os.path.join(OUT, "Sample_Input_Checklist.docx"))
    p = os.path.abspath(os.path.join(OUT, "Sample_Input_Checklist.pptx"))
    build_docx(d)
    build_pptx(p)
    print("DOCX ->", d)
    print("PPTX ->", p)
